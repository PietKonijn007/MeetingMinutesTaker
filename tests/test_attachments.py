"""Tests for the attachments feature (spec/09-attachments.md).

Covers: storage CRUD, sidecar markdown round-trip, PDF text-layer
extraction, async worker (extract + summarize), summarizer tier picking
and prompt structure, the post-append `## Attachments` helper, and the
API surface (upload/list/get/raw/delete).

The Svelte UI lands in a later batch; tests stop at the API.
"""

from __future__ import annotations

import asyncio
import io
import json
import uuid
from datetime import datetime, timezone
from pathlib import Path

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient
from sqlalchemy import create_engine, event
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

from meeting_minutes.api.deps import get_config, get_db_session
from meeting_minutes.api.routes.attachments import router
from meeting_minutes.attachments import (
    pipeline_integration as pi_mod,
)
from meeting_minutes.attachments import sidecar as sidecar_mod
from meeting_minutes.attachments import storage as storage_mod
from meeting_minutes.attachments import summarizer as summarizer_mod
from meeting_minutes.attachments import worker as worker_mod
from meeting_minutes.attachments.extractors import (
    ExtractionError,
    extract,
    extract_docx,
    extract_image,
    extract_pdf,
    extract_pptx,
)
from meeting_minutes.attachments.storage import DuplicateAttachment
from meeting_minutes.attachments.summarizer import (
    SummaryRequest,
    SummaryTier,
    build_summary_prompt,
    pick_tier,
    summarize_attachment,
)
from meeting_minutes.config import AppConfig
from meeting_minutes.models import LLMResponse
from meeting_minutes.system3.db import (
    AttachmentORM,
    MeetingORM,
    create_tables,
)


# ---------------------------------------------------------------------------
# Shared fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def session_factory():
    """In-memory SQLite that survives across threads (TestClient uses one)."""
    engine = create_engine(
        "sqlite://",
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
    )

    @event.listens_for(engine, "connect")
    def _load_vec_on_connect(dbapi_conn, connection_record):
        try:
            import sqlite_vec

            dbapi_conn.enable_load_extension(True)
            sqlite_vec.load(dbapi_conn)
        except Exception:
            pass

    create_tables(engine)
    return sessionmaker(bind=engine, autoflush=False, autocommit=False)


@pytest.fixture
def session(session_factory):
    s = session_factory()
    yield s
    s.close()


@pytest.fixture
def data_dir(tmp_path) -> Path:
    return tmp_path / "data"


@pytest.fixture
def meeting_id(session) -> str:
    """Insert a stub meeting so the FK is satisfiable."""
    mid = f"m-{uuid.uuid4().hex[:8]}"
    now = datetime.now(timezone.utc)
    session.add(
        MeetingORM(
            meeting_id=mid,
            title="Test meeting",
            date=now,
            meeting_type="standup",
            status="final",
            duration="15 minutes",
            created_at=now,
            updated_at=now,
        )
    )
    session.commit()
    return mid


def _minimal_pdf_bytes(text: str = "Hello attachment world") -> bytes:
    """Build a tiny single-page PDF using pypdf so tests don't ship a binary fixture.

    Embeds ``text`` on a single page; text-layer extraction returns it verbatim
    (modulo whitespace).
    """
    from pypdf import PdfWriter
    from pypdf.generic import (
        ArrayObject,
        DecodedStreamObject,
        DictionaryObject,
        FloatObject,
        NameObject,
        NumberObject,
    )

    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    # Build a simple content stream that uses the standard /Helvetica font.
    # Coordinates put the text near the top of the page.
    safe = text.replace("(", r"\(").replace(")", r"\)")
    content = (
        f"BT /F1 24 Tf 72 720 Td ({safe}) Tj ET".encode("latin-1")
    )
    stream = DecodedStreamObject()
    stream.set_data(content)
    page[NameObject("/Contents")] = stream

    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    resources = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font})}
    )
    page[NameObject("/Resources")] = resources
    page[NameObject("/MediaBox")] = ArrayObject(
        [NumberObject(0), NumberObject(0), FloatObject(612), FloatObject(792)]
    )

    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Sidecar
# ---------------------------------------------------------------------------


def test_sidecar_round_trip(tmp_path):
    """Write then parse — every field survives, including empty summary."""
    path = tmp_path / "abc.md"
    sidecar_mod.write_attachment_sidecar(
        path,
        frontmatter={
            "attachment_id": "abc",
            "meeting_id": "m1",
            "kind": "file",
            "title": "Q3 forecast",
            "summary_status": "pending",
        },
        extracted="--- Page 1 ---\nLine one\nLine two",
        summary="",
    )
    parsed = sidecar_mod.parse_attachment_sidecar(path)
    assert parsed.frontmatter["attachment_id"] == "abc"
    assert parsed.frontmatter["title"] == "Q3 forecast"
    assert parsed.frontmatter["summary_status"] == "pending"
    assert "Line one" in parsed.extracted
    assert parsed.summary == ""


def test_sidecar_update_summary(tmp_path):
    """update_summary preserves extracted text and sets the status fields."""
    path = tmp_path / "abc.md"
    sidecar_mod.write_attachment_sidecar(
        path,
        frontmatter={
            "attachment_id": "abc",
            "meeting_id": "m1",
            "summary_status": "pending",
        },
        extracted="EXTRACTED BODY",
        summary="",
    )

    sidecar_mod.update_summary(
        path,
        summary="Two-sentence summary.\nSecond line.",
        summary_status="ready",
        summary_target="short",
    )

    parsed = sidecar_mod.parse_attachment_sidecar(path)
    assert parsed.summary.startswith("Two-sentence summary.")
    assert parsed.extracted == "EXTRACTED BODY"
    assert parsed.frontmatter["summary_status"] == "ready"
    assert parsed.frontmatter["summary_target"] == "short"


def test_sidecar_missing_file_returns_empty(tmp_path):
    parsed = sidecar_mod.parse_attachment_sidecar(tmp_path / "does_not_exist.md")
    assert parsed.frontmatter == {}
    assert parsed.summary == ""
    assert parsed.extracted == ""


# ---------------------------------------------------------------------------
# PDF extractor
# ---------------------------------------------------------------------------


def test_extract_pdf_text_layer(tmp_path):
    pdf_bytes = _minimal_pdf_bytes("Hello attachment world")
    path = tmp_path / "tiny.pdf"
    path.write_bytes(pdf_bytes)

    text, method = extract_pdf(path)
    assert method == "pdf-text-layer"
    assert "Hello attachment world" in text
    assert "--- Page 1 ---" in text


def test_extract_dispatches_by_mime(tmp_path):
    path = tmp_path / "tiny.pdf"
    path.write_bytes(_minimal_pdf_bytes())
    text, method = extract(path, "application/pdf")
    assert method == "pdf-text-layer"
    assert text  # non-empty


def test_extract_unknown_mime_raises(tmp_path):
    path = tmp_path / "weird.xyz"
    path.write_bytes(b"whatever")
    with pytest.raises(ExtractionError):
        extract(path, "application/x-unknown")


# ---------------------------------------------------------------------------
# DOCX extractor
# ---------------------------------------------------------------------------


def _build_docx(path: Path, *, paragraphs: list[str], table_rows: list[list[str]] | None = None) -> None:
    """Helper: write a real DOCX with the given paragraphs + an optional table."""
    from docx import Document  # type: ignore

    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    if table_rows:
        table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for ri, row in enumerate(table_rows):
            for ci, cell_text in enumerate(row):
                table.rows[ri].cells[ci].text = cell_text
    doc.save(str(path))


def test_extract_docx_returns_paragraphs_in_order(tmp_path):
    path = tmp_path / "doc.docx"
    _build_docx(path, paragraphs=["First paragraph.", "Second paragraph."])

    text, method = extract_docx(path)
    assert method == "docx"
    assert text.index("First paragraph.") < text.index("Second paragraph.")


def test_extract_docx_renders_table_as_markdown(tmp_path):
    path = tmp_path / "with_table.docx"
    _build_docx(
        path,
        paragraphs=["Lead paragraph."],
        table_rows=[
            ["Name", "Owner"],
            ["Migration", "Alice"],
            ["Rollout", "Bob"],
        ],
    )

    text, method = extract_docx(path)
    assert method == "docx"
    # Table renders as `| col | col |`-style rows; verify content survived.
    assert "| Name | Owner |" in text
    assert "| Migration | Alice |" in text
    assert "| Rollout | Bob |" in text


def test_extract_docx_dispatch(tmp_path):
    path = tmp_path / "via_dispatch.docx"
    _build_docx(path, paragraphs=["dispatched"])
    text, method = extract(
        path,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    assert "dispatched" in text
    assert method == "docx"


# ---------------------------------------------------------------------------
# PPTX extractor
# ---------------------------------------------------------------------------


def _build_pptx(
    path: Path,
    *,
    slides: list[tuple[str, list[str], str | None]],  # (title, body_lines, notes)
) -> None:
    """Helper: build a real PPTX with title-only slides plus notes."""
    from pptx import Presentation  # type: ignore

    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # title-only layout
    for title, body_lines, notes in slides:
        slide = prs.slides.add_slide(blank_layout)
        if title and slide.shapes.title is not None:
            slide.shapes.title.text = title
        if body_lines:
            # Add a separate text box for body content.
            from pptx.util import Inches  # type: ignore

            tx = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(8), Inches(4)
            ).text_frame
            tx.text = body_lines[0]
            for extra in body_lines[1:]:
                tx.add_paragraph().text = extra
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    prs.save(str(path))


def test_extract_pptx_includes_titles_bodies_and_notes(tmp_path):
    path = tmp_path / "deck.pptx"
    _build_pptx(
        path,
        slides=[
            ("Q3 forecast", ["ARR landed at $12.7M", "EMEA mid-market drove growth"], "Talk about EMEA wins"),
            ("Risks", ["Renewals slipping in APAC"], None),
        ],
    )

    text, method = extract_pptx(path)
    assert method == "pptx"
    # Slide markers + titles
    assert "--- Slide 1: Q3 forecast ---" in text
    assert "--- Slide 2: Risks ---" in text
    # Body content survives
    assert "$12.7M" in text
    assert "EMEA mid-market drove growth" in text
    assert "Renewals slipping in APAC" in text
    # Speaker notes carry through with a label so the LLM knows what they are
    assert "Speaker notes:" in text
    assert "Talk about EMEA wins" in text


def test_extract_pptx_blank_slide_renders_placeholder(tmp_path):
    path = tmp_path / "thin.pptx"
    _build_pptx(path, slides=[("Title only", [], None)])

    text, _ = extract_pptx(path)
    assert "--- Slide 1: Title only ---" in text
    assert "no text content" in text


# ---------------------------------------------------------------------------
# Image OCR
# ---------------------------------------------------------------------------


def _build_image(path: Path, color=(255, 255, 255), size=(100, 50)) -> None:
    """Helper: write a tiny PNG so we have an image file on disk."""
    from PIL import Image  # type: ignore

    Image.new("RGB", size, color=color).save(str(path), format="PNG")


def test_extract_image_uses_ocr_method(tmp_path, monkeypatch):
    """Image OCR happy path with pytesseract stubbed out — no binary needed."""
    import pytesseract  # type: ignore

    path = tmp_path / "screenshot.png"
    _build_image(path)

    monkeypatch.setattr(pytesseract, "get_tesseract_version", lambda: "5.0.0")
    monkeypatch.setattr(
        pytesseract, "image_to_string", lambda img, **kw: "OCR'd line one\nLine two"
    )

    text, method = extract_image(path)
    assert method == "ocr"
    assert "OCR'd line one" in text
    assert "Line two" in text


def test_extract_image_dispatch_via_mime(tmp_path, monkeypatch):
    """Dispatcher routes image/* mime to extract_image."""
    import pytesseract  # type: ignore

    path = tmp_path / "x.png"
    _build_image(path)
    monkeypatch.setattr(pytesseract, "get_tesseract_version", lambda: "5.0.0")
    monkeypatch.setattr(pytesseract, "image_to_string", lambda img, **kw: "ok")

    text, method = extract(path, "image/png")
    assert method == "ocr"
    assert text == "ok"


def test_extract_image_missing_binary_raises_with_install_hint(tmp_path, monkeypatch):
    """When tesseract isn't on PATH, surface a clear error — not a stack trace."""
    import pytesseract  # type: ignore

    path = tmp_path / "x.png"
    _build_image(path)

    def _boom():
        raise EnvironmentError("tesseract is not installed")

    monkeypatch.setattr(pytesseract, "get_tesseract_version", _boom)

    with pytest.raises(ExtractionError) as exc_info:
        extract_image(path)
    assert "tesseract" in str(exc_info.value).lower()
    assert "brew install tesseract" in str(exc_info.value)


def test_extract_image_empty_result_is_not_an_error(tmp_path, monkeypatch):
    """A blank image OCR'ing to nothing is valid — return empty string, no raise."""
    import pytesseract  # type: ignore

    path = tmp_path / "blank.png"
    _build_image(path)

    monkeypatch.setattr(pytesseract, "get_tesseract_version", lambda: "5.0.0")
    monkeypatch.setattr(pytesseract, "image_to_string", lambda img, **kw: "")

    text, method = extract_image(path)
    assert text == ""
    assert method == "ocr"


def test_extract_image_ocr_failure_wraps_to_extraction_error(tmp_path, monkeypatch):
    import pytesseract  # type: ignore

    path = tmp_path / "x.png"
    _build_image(path)
    monkeypatch.setattr(pytesseract, "get_tesseract_version", lambda: "5.0.0")

    def _raise(img, **kw):
        raise RuntimeError("OCR engine crashed")

    monkeypatch.setattr(pytesseract, "image_to_string", _raise)
    with pytest.raises(ExtractionError) as exc_info:
        extract_image(path)
    assert "OCR failed" in str(exc_info.value)


# ---------------------------------------------------------------------------
# Storage CRUD
# ---------------------------------------------------------------------------


def test_add_file_persists_row_and_disk(session, data_dir, meeting_id):
    raw = b"some bytes"
    row = storage_mod.add_file(
        session=session,
        data_dir=data_dir,
        meeting_id=meeting_id,
        fileobj=io.BytesIO(raw),
        original_filename="hello.pdf",
        mime_type="application/pdf",
        title="Hello",
    )
    session.commit()

    assert row.attachment_id
    assert row.kind == "file"
    assert row.size_bytes == len(raw)
    assert row.status == "pending"
    assert row.title == "Hello"

    on_disk = storage_mod.original_path(
        data_dir, meeting_id, row.attachment_id, ".pdf"
    )
    assert on_disk.exists()
    assert on_disk.read_bytes() == raw


def test_add_file_dedupe_returns_existing(session, data_dir, meeting_id):
    raw = b"identical content"
    first = storage_mod.add_file(
        session=session,
        data_dir=data_dir,
        meeting_id=meeting_id,
        fileobj=io.BytesIO(raw),
        original_filename="a.pdf",
        mime_type="application/pdf",
    )
    session.commit()

    with pytest.raises(DuplicateAttachment) as exc_info:
        storage_mod.add_file(
            session=session,
            data_dir=data_dir,
            meeting_id=meeting_id,
            fileobj=io.BytesIO(raw),
            original_filename="b.pdf",
            mime_type="application/pdf",
        )
    assert exc_info.value.attachment_id == first.attachment_id


def test_add_file_unknown_meeting_raises(session, data_dir):
    with pytest.raises(ValueError):
        storage_mod.add_file(
            session=session,
            data_dir=data_dir,
            meeting_id="does-not-exist",
            fileobj=io.BytesIO(b"x"),
            original_filename="x.pdf",
            mime_type="application/pdf",
        )


def test_list_for_meeting_orders_by_created(session, data_dir, meeting_id):
    for i in range(3):
        storage_mod.add_file(
            session=session,
            data_dir=data_dir,
            meeting_id=meeting_id,
            fileobj=io.BytesIO(f"content {i}".encode()),
            original_filename=f"file{i}.pdf",
            mime_type="application/pdf",
        )
    session.commit()
    rows = storage_mod.list_for_meeting(session, meeting_id)
    assert len(rows) == 3
    assert rows[0].original_filename == "file0.pdf"
    assert rows[2].original_filename == "file2.pdf"


def test_delete_drops_row_and_files(session, data_dir, meeting_id, tmp_path):
    row = storage_mod.add_file(
        session=session,
        data_dir=data_dir,
        meeting_id=meeting_id,
        fileobj=io.BytesIO(b"goodbye"),
        original_filename="bye.pdf",
        mime_type="application/pdf",
    )
    session.commit()
    aid = row.attachment_id
    on_disk = storage_mod.original_path(data_dir, meeting_id, aid, ".pdf")
    sidecar_path = storage_mod.sidecar_path(data_dir, meeting_id, aid)
    sidecar_path.parent.mkdir(parents=True, exist_ok=True)
    sidecar_path.write_text("placeholder", encoding="utf-8")

    assert storage_mod.delete(session, data_dir, aid) is True
    session.commit()
    assert session.get(AttachmentORM, aid) is None
    assert not on_disk.exists()
    assert not sidecar_path.exists()


def test_delete_unknown_returns_false(session, data_dir):
    assert storage_mod.delete(session, data_dir, "no-such-id") is False


# ---------------------------------------------------------------------------
# Worker
# ---------------------------------------------------------------------------


def test_worker_extracts_pdf_and_writes_sidecar(
    session_factory, data_dir, meeting_id
):
    """End-to-end: row + file on disk → worker → sidecar populated, status=ready.

    Uses a fake LLM so the worker's summary phase doesn't try to call out
    to a real provider. The full extract+summarize round-trip is exercised
    in :func:`test_worker_runs_extraction_and_summary`; this test focuses
    on the extraction half producing a well-formed sidecar.
    """
    session = session_factory()
    try:
        row = storage_mod.add_file(
            session=session,
            data_dir=data_dir,
            meeting_id=meeting_id,
            fileobj=io.BytesIO(_minimal_pdf_bytes("Pipeline-test text")),
            original_filename="doc.pdf",
            mime_type="application/pdf",
        )
        session.commit()
        attachment_id = row.attachment_id
    finally:
        session.close()

    cfg = AppConfig(data_dir=str(data_dir))

    asyncio.run(
        worker_mod.process_attachment(
            cfg, attachment_id,
            session_factory=session_factory,
            llm_client=_FakeLLM(response_text="placeholder summary"),
        )
    )

    # Re-read from a fresh session to see committed state.
    session = session_factory()
    try:
        refreshed = session.get(AttachmentORM, attachment_id)
        assert refreshed is not None
        assert refreshed.status == "ready"

        sidecar = sidecar_mod.parse_attachment_sidecar(
            storage_mod.sidecar_path(data_dir, meeting_id, attachment_id)
        )
        assert sidecar.frontmatter["extraction_method"] == "pdf-text-layer"
        assert "Pipeline-test text" in sidecar.extracted
    finally:
        session.close()


def test_worker_records_error_when_file_missing(
    session_factory, data_dir, meeting_id
):
    """If the original file is gone, worker flips status=error cleanly."""
    session = session_factory()
    try:
        row = storage_mod.add_file(
            session=session,
            data_dir=data_dir,
            meeting_id=meeting_id,
            fileobj=io.BytesIO(b"placeholder"),
            original_filename="x.pdf",
            mime_type="application/pdf",
        )
        session.commit()
        attachment_id = row.attachment_id
        # Delete the on-disk original out from under the worker.
        storage_mod.original_path(data_dir, meeting_id, attachment_id, ".pdf").unlink()
    finally:
        session.close()

    cfg = AppConfig(data_dir=str(data_dir))
    asyncio.run(
        worker_mod.process_attachment(
            cfg, attachment_id, session_factory=session_factory
        )
    )

    session = session_factory()
    try:
        refreshed = session.get(AttachmentORM, attachment_id)
        assert refreshed.status == "error"
        assert refreshed.error and "missing" in refreshed.error.lower()
    finally:
        session.close()


# ---------------------------------------------------------------------------
# API
# ---------------------------------------------------------------------------


@pytest.fixture
def app(session_factory, data_dir):
    """A FastAPI app with only the attachments router mounted."""
    app = FastAPI()
    app.include_router(router)

    def override_session():
        s = session_factory()
        try:
            yield s
        finally:
            s.close()

    def override_config():
        return AppConfig(data_dir=str(data_dir))

    app.dependency_overrides[get_db_session] = override_session
    app.dependency_overrides[get_config] = override_config
    return app


@pytest.fixture
def client(app, monkeypatch):
    """TestClient with the worker stubbed out — handler tests don't need to
    actually run extraction; that's covered separately."""

    def _no_op_schedule(_config, _attachment_id):
        async def _noop():
            return None

        return asyncio.get_event_loop().create_task(_noop())

    monkeypatch.setattr(worker_mod, "schedule", _no_op_schedule)
    return TestClient(app)


def test_api_upload_creates_row(client, meeting_id, data_dir):
    pdf = _minimal_pdf_bytes("API upload body")
    response = client.post(
        f"/api/meetings/{meeting_id}/attachments",
        files={"file": ("hello.pdf", pdf, "application/pdf")},
        data={"title": "Hello", "caption": "from a test"},
    )
    assert response.status_code == 202, response.text
    payload = response.json()
    assert payload["title"] == "Hello"
    assert payload["caption"] == "from a test"
    assert payload["mime_type"] == "application/pdf"
    assert payload["status"] == "pending"

    # File ended up on disk under the meeting's folder.
    target_dir = data_dir / "attachments" / meeting_id
    assert target_dir.exists()
    files = list(target_dir.iterdir())
    assert any(f.suffix == ".pdf" for f in files)


def test_api_upload_unknown_meeting_404(client):
    response = client.post(
        "/api/meetings/no-such-meeting/attachments",
        files={"file": ("x.pdf", b"data", "application/pdf")},
    )
    assert response.status_code == 404


def test_api_upload_rejects_disallowed_mime(client, meeting_id):
    response = client.post(
        f"/api/meetings/{meeting_id}/attachments",
        files={"file": ("evil.exe", b"MZ\x00\x00", "application/x-msdownload")},
    )
    assert response.status_code == 415


def test_api_upload_dedup_returns_existing(client, meeting_id):
    pdf = _minimal_pdf_bytes("dedup body")
    first = client.post(
        f"/api/meetings/{meeting_id}/attachments",
        files={"file": ("a.pdf", pdf, "application/pdf")},
    )
    assert first.status_code == 202
    aid = first.json()["attachment_id"]

    second = client.post(
        f"/api/meetings/{meeting_id}/attachments",
        files={"file": ("b.pdf", pdf, "application/pdf")},
    )
    # Dedup path returns the existing summary — handler returns 202 default
    # for the route, so we accept either 200 or 202 as long as the id matches.
    assert second.status_code in (200, 202)
    assert second.json()["attachment_id"] == aid


def test_api_list_returns_uploads(client, meeting_id):
    for i in range(2):
        client.post(
            f"/api/meetings/{meeting_id}/attachments",
            files={"file": (f"f{i}.pdf", _minimal_pdf_bytes(f"body {i}"), "application/pdf")},
        )

    response = client.get(f"/api/meetings/{meeting_id}/attachments")
    assert response.status_code == 200
    payload = response.json()
    assert len(payload) == 2
    filenames = sorted(item["original_filename"] for item in payload)
    assert filenames == ["f0.pdf", "f1.pdf"]


def test_api_get_detail_includes_sidecar(
    client, session_factory, data_dir, meeting_id
):
    """If a sidecar exists, the detail endpoint surfaces it."""
    upload = client.post(
        f"/api/meetings/{meeting_id}/attachments",
        files={"file": ("doc.pdf", _minimal_pdf_bytes("body"), "application/pdf")},
    )
    aid = upload.json()["attachment_id"]

    # Stub sidecar so we don't need the worker to actually run.
    sidecar_mod.write_attachment_sidecar(
        storage_mod.sidecar_path(data_dir, meeting_id, aid),
        frontmatter={
            "attachment_id": aid,
            "meeting_id": meeting_id,
            "summary_status": "ready",
        },
        extracted="EXTRACTED",
        summary="Short summary.",
    )

    response = client.get(f"/api/attachments/{aid}")
    assert response.status_code == 200
    payload = response.json()
    assert payload["summary"] == "Short summary."
    assert payload["extracted_text"] == "EXTRACTED"
    assert payload["summary_status"] == "ready"


def test_api_raw_returns_original_bytes(client, meeting_id):
    pdf = _minimal_pdf_bytes("raw body")
    upload = client.post(
        f"/api/meetings/{meeting_id}/attachments",
        files={"file": ("doc.pdf", pdf, "application/pdf")},
    )
    aid = upload.json()["attachment_id"]

    response = client.get(f"/api/attachments/{aid}/raw")
    assert response.status_code == 200
    assert response.headers["content-type"].startswith("application/pdf")
    assert response.content == pdf


def test_api_delete_removes_row(client, session_factory, data_dir, meeting_id):
    upload = client.post(
        f"/api/meetings/{meeting_id}/attachments",
        files={"file": ("doc.pdf", _minimal_pdf_bytes("body"), "application/pdf")},
    )
    aid = upload.json()["attachment_id"]

    response = client.delete(f"/api/attachments/{aid}")
    assert response.status_code == 204

    # Subsequent GET 404s.
    follow = client.get(f"/api/attachments/{aid}")
    assert follow.status_code == 404


# ---------------------------------------------------------------------------
# Summarizer (tier picking + prompt structure)
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "size, expected",
    [
        (0, SummaryTier.SHORT),
        (200, SummaryTier.SHORT),
        (499, SummaryTier.SHORT),
        (500, SummaryTier.MEDIUM),
        (4_999, SummaryTier.MEDIUM),
        (5_000, SummaryTier.LONG),
        (29_999, SummaryTier.LONG),
        (30_000, SummaryTier.XLONG),
        (200_000, SummaryTier.XLONG),
    ],
)
def test_summary_tier_thresholds(size, expected):
    text = "x" * size
    assert pick_tier(text) == expected


def test_build_summary_prompt_short_includes_caption_and_tier_instruction():
    req = SummaryRequest(
        title="One slide",
        caption="Q3 ARR slide we discussed at minute 12",
        source="q3.pdf",
        extraction_method="pdf-text-layer",
        extracted_text="ARR: $12.7M",
    )
    system, user, tier, truncated = build_summary_prompt(req)
    assert tier == SummaryTier.SHORT
    assert truncated is False
    assert "CRITICAL RULES" in system
    assert "Q3 ARR slide" in user  # caption surfaces
    assert "CONCISE" in user  # short-tier instruction
    assert "$12.7M" in user
    # Source + title rendered explicitly so the LLM sees the framing.
    assert "Source: q3.pdf" in user
    assert "Title: One slide" in user


def test_build_summary_prompt_long_uses_subheaders_instruction():
    req = SummaryRequest(
        title="Big PDF",
        caption=None,
        source="big.pdf",
        extraction_method="pdf-text-layer",
        extracted_text="x" * 6_000,
    )
    _, user, tier, _ = build_summary_prompt(req)
    assert tier == SummaryTier.LONG
    assert "DETAILED" in user
    assert "###" in user  # subheader guidance


def test_build_summary_prompt_truncates_oversized_input():
    big = "x" * 200_000
    req = SummaryRequest(
        title="Huge", caption=None, source="huge.pdf",
        extraction_method="pdf-text-layer", extracted_text=big,
    )
    _, user, tier, truncated = build_summary_prompt(req)
    assert tier == SummaryTier.XLONG
    assert truncated is True
    assert "[... document truncated ...]" in user
    # We don't blow the prompt up to 200k chars — the tail is replaced.
    assert "[... document truncated ...]" in user
    assert user.count("x") < 110_000


def test_build_summary_prompt_ocr_marks_skepticism():
    req = SummaryRequest(
        title="Whiteboard photo", caption=None, source="board.png",
        extraction_method="ocr", extracted_text="meeting agenda",
    )
    _, user, _, _ = build_summary_prompt(req)
    assert "OCR output more skeptically" in user


class _FakeLLM:
    """Minimal LLMClient stand-in for summarizer tests.

    Records the calls it received so tests can assert on the prompts the
    summarizer constructed without standing up a real LLM.
    """

    def __init__(self, response_text: str = "Summary body."):
        self.response_text = response_text
        self.calls: list[tuple[str, str]] = []  # (system_prompt, prompt)

    async def generate(self, prompt: str, system_prompt: str = "") -> LLMResponse:
        self.calls.append((system_prompt, prompt))
        return LLMResponse(
            text=self.response_text,
            provider="fake",
            model="fake",
            input_tokens=100,
            output_tokens=50,
            cost_usd=0.0,
            processing_time_seconds=0.01,
        )


def test_summarize_attachment_calls_llm_and_returns_markdown():
    fake = _FakeLLM(response_text="### Bullet\n- one\n- two\n")
    req = SummaryRequest(
        title="t", caption="c", source="s", extraction_method="pdf-text-layer",
        extracted_text="some content here",
    )
    result = asyncio.run(summarize_attachment(fake, req))
    assert result.summary_markdown.startswith("### Bullet")
    assert result.tier == SummaryTier.SHORT
    assert result.truncated is False
    assert len(fake.calls) == 1


def test_summarize_attachment_skips_llm_for_empty_extraction():
    """Calling the LLM with no content wastes tokens; short-circuit instead."""
    fake = _FakeLLM()
    req = SummaryRequest(
        title="empty", caption=None, source="x", extraction_method="ocr",
        extracted_text="",
    )
    result = asyncio.run(summarize_attachment(fake, req))
    assert "Could not extract text" in result.summary_markdown
    assert fake.calls == []


# ---------------------------------------------------------------------------
# Worker — full extract + summarize round trip
# ---------------------------------------------------------------------------


def test_worker_runs_extraction_and_summary(
    session_factory, data_dir, meeting_id
):
    """Full worker path: extract → summarize → sidecar populated, status=ready."""
    session = session_factory()
    try:
        row = storage_mod.add_file(
            session=session,
            data_dir=data_dir,
            meeting_id=meeting_id,
            fileobj=io.BytesIO(_minimal_pdf_bytes("Pipeline integration body")),
            original_filename="doc.pdf",
            mime_type="application/pdf",
            title="Integration doc",
            caption="why this matters",
        )
        session.commit()
        attachment_id = row.attachment_id
    finally:
        session.close()

    fake_llm = _FakeLLM(response_text="A grounded summary mentioning $12.7M.")
    cfg = AppConfig(data_dir=str(data_dir))
    asyncio.run(
        worker_mod.process_attachment(
            cfg, attachment_id,
            session_factory=session_factory,
            llm_client=fake_llm,
        )
    )

    session = session_factory()
    try:
        refreshed = session.get(AttachmentORM, attachment_id)
        assert refreshed.status == "ready"

        sidecar = sidecar_mod.parse_attachment_sidecar(
            storage_mod.sidecar_path(data_dir, meeting_id, attachment_id)
        )
        assert sidecar.frontmatter["summary_status"] == "ready"
        assert sidecar.frontmatter["summary_target"] == SummaryTier.SHORT.value
        assert "$12.7M" in sidecar.summary
        assert "Pipeline integration body" in sidecar.extracted
    finally:
        session.close()
    # The summarizer was actually called once; the prompt carried
    # the title and caption from the row.
    assert len(fake_llm.calls) == 1
    _, user_prompt = fake_llm.calls[0]
    assert "Integration doc" in user_prompt
    assert "why this matters" in user_prompt


def test_worker_records_summary_error_but_keeps_row_ready(
    session_factory, data_dir, meeting_id
):
    """Extraction succeeded, summary failed: row is ready, sidecar marks error."""

    class _BoomLLM(_FakeLLM):
        async def generate(self, prompt: str, system_prompt: str = "") -> LLMResponse:  # type: ignore[override]
            raise RuntimeError("LLM provider unreachable")

    session = session_factory()
    try:
        row = storage_mod.add_file(
            session=session,
            data_dir=data_dir,
            meeting_id=meeting_id,
            fileobj=io.BytesIO(_minimal_pdf_bytes("body")),
            original_filename="doc.pdf",
            mime_type="application/pdf",
        )
        session.commit()
        attachment_id = row.attachment_id
    finally:
        session.close()

    cfg = AppConfig(data_dir=str(data_dir))
    asyncio.run(
        worker_mod.process_attachment(
            cfg, attachment_id,
            session_factory=session_factory,
            llm_client=_BoomLLM(),
        )
    )

    session = session_factory()
    try:
        refreshed = session.get(AttachmentORM, attachment_id)
        # Extraction succeeded → row stays ready (no point re-extracting
        # to retry the summary). The summary error is on the sidecar
        # where the UI picks it up.
        assert refreshed.status == "ready"
    finally:
        session.close()

    sidecar = sidecar_mod.parse_attachment_sidecar(
        storage_mod.sidecar_path(data_dir, meeting_id, attachment_id)
    )
    assert sidecar.frontmatter["summary_status"] == "error"
    assert "LLM provider unreachable" in sidecar.frontmatter.get("summary_error", "")


# ---------------------------------------------------------------------------
# Pipeline integration — gather, wait, render, post-append
# ---------------------------------------------------------------------------


def _seed_sidecar(
    data_dir: Path,
    meeting_id: str,
    *,
    attachment_id: str,
    title: str,
    summary_status: str = "ready",
    summary: str = "Some summary.",
    extracted: str = "Some extracted text.",
    source: str = "doc.pdf",
):
    """Helper: write a sidecar file in the meeting's attachments folder."""
    path = storage_mod.sidecar_path(data_dir, meeting_id, attachment_id)
    sidecar_mod.write_attachment_sidecar(
        path,
        frontmatter={
            "attachment_id": attachment_id,
            "meeting_id": meeting_id,
            "kind": "file",
            "title": title,
            "source": source,
            "extraction_method": "pdf-text-layer",
            "summary_status": summary_status,
        },
        extracted=extracted,
        summary=summary,
    )
    return path


def test_gather_entries_returns_empty_when_no_folder(data_dir):
    assert pi_mod.gather_entries(data_dir, "no-such-meeting") == []


def test_gather_entries_reads_sidecars_in_filename_order(data_dir):
    mid = "m-1"
    _seed_sidecar(data_dir, mid, attachment_id="a-aaa", title="First")
    _seed_sidecar(data_dir, mid, attachment_id="a-zzz", title="Last")
    entries = pi_mod.gather_entries(data_dir, mid)
    assert [e.attachment_id for e in entries] == ["a-aaa", "a-zzz"]
    assert entries[0].title == "First"


def test_render_llm_context_block_skips_pending_and_errored(data_dir):
    entries = [
        pi_mod.AttachmentEntry(
            attachment_id="a1", title="Ready", source="r.pdf",
            summary="Body content.", summary_status="ready",
            extraction_method="pdf-text-layer", kind="file",
        ),
        pi_mod.AttachmentEntry(
            attachment_id="a2", title="Pending", source="p.pdf",
            summary="", summary_status="pending",
            extraction_method="pdf-text-layer", kind="file",
        ),
        pi_mod.AttachmentEntry(
            attachment_id="a3", title="Errored", source="e.pdf",
            summary="", summary_status="error",
            extraction_method="pdf-text-layer", kind="file",
        ),
    ]
    block = pi_mod.render_llm_context_block(entries)
    assert "ATTACHED MATERIAL: Ready" in block
    assert "Body content." in block
    assert "Pending" not in block
    assert "Errored" not in block
    # Preamble explains the contract to the LLM.
    assert "ground-truth context" in block


def test_render_llm_context_block_empty_when_nothing_ready():
    assert pi_mod.render_llm_context_block([]) == ""


def test_wait_for_pending_returns_when_none_pending(data_dir):
    mid = "m-1"
    _seed_sidecar(data_dir, mid, attachment_id="a1", title="One",
                  summary_status="ready")
    entries = asyncio.run(
        pi_mod.wait_for_pending(data_dir, mid, timeout_s=2.0)
    )
    assert len(entries) == 1


def test_wait_for_pending_times_out(data_dir):
    """A persistently-pending sidecar means the pipeline can't wait forever."""
    mid = "m-1"
    _seed_sidecar(data_dir, mid, attachment_id="a1", title="Stuck",
                  summary_status="pending")
    entries = asyncio.run(
        pi_mod.wait_for_pending(data_dir, mid, timeout_s=0.2,
                                poll_interval_s=0.05)
    )
    # Returns whatever was on disk at the deadline; doesn't raise.
    assert len(entries) == 1
    assert entries[0].summary_status == "pending"


def test_append_attachments_section_writes_md_and_json(tmp_path):
    """Updates both .md and .json so DB ingestion picks up the appended block."""
    minutes_dir = tmp_path / "data" / "minutes"
    minutes_dir.mkdir(parents=True)
    md_path = minutes_dir / "m-1.md"
    json_path = minutes_dir / "m-1.json"
    md_path.write_text("# Title\n\n## Summary\nbody\n", encoding="utf-8")
    json_path.write_text(
        '{"meeting_id": "m-1", "minutes_markdown": "# Title\\n\\n## Summary\\nbody\\n"}',
        encoding="utf-8",
    )

    entries = [
        pi_mod.AttachmentEntry(
            attachment_id="a1", title="Slide deck", source="deck.pdf",
            summary="Key takeaways: $12.7M ARR.", summary_status="ready",
            extraction_method="pdf-text-layer", kind="file",
        ),
    ]
    final = pi_mod.append_attachments_section_to_files(
        tmp_path / "data", "m-1", entries,
    )
    assert final is not None
    assert "## Attachments" in final
    assert "Slide deck" in final
    assert "$12.7M ARR" in final
    assert "[View source](/api/attachments/a1/raw)" in final

    # Both files reflect the change.
    assert "## Attachments" in md_path.read_text(encoding="utf-8")
    json_data = json.loads(json_path.read_text(encoding="utf-8"))
    assert "## Attachments" in json_data["minutes_markdown"]


def test_append_attachments_section_idempotent(tmp_path):
    """Two consecutive calls converge — no duplicated section."""
    minutes_dir = tmp_path / "data" / "minutes"
    minutes_dir.mkdir(parents=True)
    md_path = minutes_dir / "m-1.md"
    md_path.write_text("# Title\n\nbody\n", encoding="utf-8")
    entries = [
        pi_mod.AttachmentEntry(
            attachment_id="a1", title="Doc", source="d.pdf",
            summary="Body 1.", summary_status="ready",
            extraction_method="pdf-text-layer", kind="file",
        ),
    ]
    pi_mod.append_attachments_section_to_files(tmp_path / "data", "m-1", entries)
    pi_mod.append_attachments_section_to_files(tmp_path / "data", "m-1", entries)

    md = md_path.read_text(encoding="utf-8")
    assert md.count("## Attachments") == 1
    assert md.count("Body 1.") == 1


def test_append_attachments_section_replaces_stale_block_after_change(tmp_path):
    """When the attachments list changes, the rendered section reflects only the new state."""
    minutes_dir = tmp_path / "data" / "minutes"
    minutes_dir.mkdir(parents=True)
    md_path = minutes_dir / "m-1.md"
    md_path.write_text("# Title\n\nbody\n", encoding="utf-8")

    initial = [
        pi_mod.AttachmentEntry(
            attachment_id="a1", title="Old", source="o.pdf",
            summary="Old summary.", summary_status="ready",
            extraction_method="pdf-text-layer", kind="file",
        ),
    ]
    pi_mod.append_attachments_section_to_files(tmp_path / "data", "m-1", initial)
    assert "Old summary." in md_path.read_text(encoding="utf-8")

    updated = [
        pi_mod.AttachmentEntry(
            attachment_id="a2", title="New", source="n.pdf",
            summary="New summary.", summary_status="ready",
            extraction_method="pdf-text-layer", kind="file",
        ),
    ]
    pi_mod.append_attachments_section_to_files(tmp_path / "data", "m-1", updated)
    md = md_path.read_text(encoding="utf-8")
    assert "Old summary." not in md
    assert "New summary." in md


def test_append_attachments_section_handles_errored_summary(tmp_path):
    """Errored summary still gets a row in the rendered section, with a clear note."""
    minutes_dir = tmp_path / "data" / "minutes"
    minutes_dir.mkdir(parents=True)
    md_path = minutes_dir / "m-1.md"
    md_path.write_text("# Title\n\nbody\n", encoding="utf-8")

    entries = [
        pi_mod.AttachmentEntry(
            attachment_id="a1", title="Failed", source="f.pdf",
            summary="", summary_status="error",
            extraction_method="ocr", kind="image",
        ),
    ]
    pi_mod.append_attachments_section_to_files(tmp_path / "data", "m-1", entries)
    md = md_path.read_text(encoding="utf-8")
    assert "Failed" in md
    assert "Summary failed" in md


def test_append_attachments_section_empty_strips_stale_block(tmp_path):
    """Calling with no entries removes any prior section but doesn't add a new one."""
    minutes_dir = tmp_path / "data" / "minutes"
    minutes_dir.mkdir(parents=True)
    md_path = minutes_dir / "m-1.md"
    md_path.write_text(
        "# Title\n\nbody\n\n## Attachments\n\n### Stale\nOld.\n",
        encoding="utf-8",
    )
    final = pi_mod.append_attachments_section_to_files(
        tmp_path / "data", "m-1", []
    )
    assert final is None
    md = md_path.read_text(encoding="utf-8")
    assert "## Attachments" not in md
    assert "Stale" not in md
    assert "body" in md  # rest of the file untouched
