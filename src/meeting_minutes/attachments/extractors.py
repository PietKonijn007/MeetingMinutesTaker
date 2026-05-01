"""Per-kind text extraction.

Covers PDF text-layer, image OCR, DOCX, and PPTX. Scanned-PDF OCR
fallback (rendering pages via Poppler then OCRing them) and link-fetch
land in a later batch — see ``specs/09-attachments.md`` for the table.

Each extractor returns ``(extracted_text, extraction_method)``. The
``extraction_method`` string is recorded on the sidecar frontmatter so
future readers can tell *how* the text came to be — used by the
summarizer prompt to mark OCR output as skeptical, and useful for
debugging when a summary is unexpectedly thin.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Callable

_LOG = logging.getLogger(__name__)

# MIME types we explicitly recognize. Anything outside this set falls
# through to the extension dispatcher; the upload handler enforces a
# stricter allowlist already, so we don't need to police that here.
_PDF_MIME = "application/pdf"
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

# Image extensions we OCR. HEIC requires pillow-heif and ships in a
# follow-up batch; for now an HEIC upload will reach the dispatcher and
# raise a clear ExtractionError.
_IMAGE_EXTS = {".png", ".jpg", ".jpeg"}


class ExtractionError(RuntimeError):
    """Raised when extraction fails for any reason. Worker turns this into
    a sidecar ``summary_status: error`` and a DB ``status='error'``.
    """


def extract(path: Path, mime_type: str | None) -> tuple[str, str]:
    """Dispatch to the right extractor by mime type.

    Falls back on file extension when ``mime_type`` is missing or
    unrecognized — clients sometimes upload PDFs with a generic
    ``application/octet-stream``.
    """
    handler = _resolve_handler(path, mime_type)
    if handler is None:
        raise ExtractionError(
            f"No extractor for mime={mime_type!r} ext={path.suffix!r}"
        )
    return handler(path)


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------


def _resolve_handler(
    path: Path, mime_type: str | None
) -> Callable[[Path], tuple[str, str]] | None:
    ext = path.suffix.lower()
    if mime_type == _PDF_MIME or ext == ".pdf":
        return extract_pdf
    if mime_type == _DOCX_MIME or ext == ".docx":
        return extract_docx
    if mime_type == _PPTX_MIME or ext == ".pptx":
        return extract_pptx
    if (mime_type and mime_type.startswith("image/")) or ext in _IMAGE_EXTS:
        return extract_image
    return None


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def extract_pdf(path: Path) -> tuple[str, str]:
    """Extract text from a PDF's text layer.

    Returns ``(text, "pdf-text-layer")``. Pages are joined with a clear
    ``--- Page N ---`` separator so the summarizer can refer to specific
    pages later.

    If the text layer is empty or near-empty (a scanned PDF), the OCR
    fallback in the next batch will take over. For now we surface what
    pypdf gives us — possibly empty — and log a warning. Empty extractions
    still produce a valid sidecar; the summarizer will note "no text
    extracted" when it runs.
    """
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ExtractionError(
            "pypdf is not installed. Run `pip install pypdf>=4.0.0` "
            "or reinstall the project to pick up the new dependency."
        ) from exc

    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        raise ExtractionError(f"Could not open PDF: {exc}") from exc

    if reader.is_encrypted:
        # Try the empty password — some PDFs are "encrypted" with no
        # password, which pypdf surfaces as is_encrypted=True until you
        # decrypt with "". If that fails, give up cleanly.
        try:
            reader.decrypt("")
        except Exception:
            raise ExtractionError("PDF is password-protected; cannot extract.")

    pages: list[str] = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception as exc:
            _LOG.warning("PDF %s: page %d extraction failed: %s", path.name, idx, exc)
            page_text = ""
        if page_text.strip():
            pages.append(f"--- Page {idx} ---\n{page_text.strip()}")

    body = "\n\n".join(pages)
    if not body.strip():
        _LOG.info(
            "PDF %s: text layer is empty — likely scanned, awaiting OCR fallback.",
            path.name,
        )
    return body, "pdf-text-layer"


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def extract_docx(path: Path) -> tuple[str, str]:
    """Extract text from a DOCX in document order.

    Walks paragraphs and tables top-to-bottom so the summarizer sees
    content in the same flow the author wrote it. Tables get rendered
    as ``| col1 | col2 |``-style markdown rows — pleasant to read in
    the sidecar and carries enough structure for the LLM to recognize
    they were a table.

    ``python-docx`` is a direct dep, so an ImportError here is a real
    install problem. Surface it as ExtractionError with a clear message.
    """
    try:
        from docx import Document  # type: ignore
    except ImportError as exc:
        raise ExtractionError(
            "python-docx is not installed. Reinstall the project to pick "
            "up the new dependency."
        ) from exc

    try:
        doc = Document(str(path))
    except Exception as exc:
        raise ExtractionError(f"Could not open DOCX: {exc}") from exc

    parts: list[str] = []
    body = doc.element.body  # underlying lxml element preserves doc order

    # python-docx doesn't expose paragraphs+tables in document order
    # via its high-level API — walking the underlying XML is the
    # supported workaround. Each w:p / w:tbl child we hit gets mapped
    # back to its high-level wrapper.
    from docx.oxml.ns import qn  # type: ignore
    from docx.table import Table  # type: ignore
    from docx.text.paragraph import Paragraph  # type: ignore

    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            text = Paragraph(child, doc).text
            if text.strip():
                parts.append(text)
        elif child.tag == qn("w:tbl"):
            parts.append(_render_docx_table(Table(child, doc)))

    body_text = "\n\n".join(parts)
    if not body_text.strip():
        _LOG.info("DOCX %s: no text extracted (empty document?)", path.name)
    return body_text, "docx"


def _render_docx_table(table) -> str:
    """Render a python-docx ``Table`` as a simple markdown table.

    No header detection — DOCX tables don't necessarily have a header
    row, and we'd rather render every row honestly than guess wrong and
    hide content. Cells with line breaks are flattened to a single line
    so the table renders cleanly.
    """
    rows: list[str] = []
    for row in table.rows:
        cells = [
            " ".join(cell.text.split())  # collapse whitespace
            for cell in row.cells
        ]
        rows.append("| " + " | ".join(cells) + " |")
    return "\n".join(rows) if rows else ""


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def extract_pptx(path: Path) -> tuple[str, str]:
    """Extract text from a PPTX, slide by slide.

    Each slide is rendered as ``--- Slide N: <title> ---`` followed by
    body shapes and speaker notes. Speaker notes are critical context
    that the summarizer would otherwise miss — they're often where the
    presenter parked the actual narrative for a slide that visually
    contains only a chart or a one-liner.
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise ExtractionError(
            "python-pptx is not installed. Reinstall the project to pick "
            "up the new dependency."
        ) from exc

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise ExtractionError(f"Could not open PPTX: {exc}") from exc

    blocks: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = _pptx_slide_title(slide)
        header = f"--- Slide {idx}" + (f": {title}" if title else "") + " ---"
        body_lines: list[str] = []

        for shape in slide.shapes:
            # ``has_text_frame`` is the safe accessor — placeholders,
            # text boxes, and grouped shapes all expose it. We skip the
            # title shape (already rendered in the header) to avoid
            # duplicate output.
            if not getattr(shape, "has_text_frame", False):
                continue
            if shape == slide.shapes.title:
                continue
            text = shape.text_frame.text.strip()
            if text:
                body_lines.append(text)

        notes_text = ""
        if slide.has_notes_slide:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        if notes_text:
            body_lines.append(f"_Speaker notes:_\n{notes_text}")

        if body_lines:
            blocks.append(header + "\n" + "\n\n".join(body_lines))
        else:
            blocks.append(header + "\n_(no text content)_")

    body = "\n\n".join(blocks)
    if not body.strip():
        _LOG.info("PPTX %s: no text content found", path.name)
    return body, "pptx"


def _pptx_slide_title(slide) -> str:
    """Return the slide's title text, or ``""`` if no title placeholder."""
    title_shape = slide.shapes.title
    if title_shape is None:
        return ""
    try:
        return (title_shape.text_frame.text or "").strip()
    except Exception:
        return ""


# ---------------------------------------------------------------------------
# Image OCR
# ---------------------------------------------------------------------------


def extract_image(path: Path) -> tuple[str, str]:
    """OCR an image file via pytesseract.

    Requires the ``tesseract`` binary on PATH (the Python wrapper shells
    out to it). When the binary is missing we surface a clear install
    hint so the user can fix it themselves rather than seeing a generic
    "command not found".

    Returns ``("", "ocr")`` for images that produce no text — a
    decorative screenshot is still a valid attachment, the summarizer
    just gets nothing to ground on. Empty results are fine.
    """
    try:
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore
    except ImportError as exc:
        raise ExtractionError(
            "pytesseract or Pillow is not installed. Reinstall the project "
            "to pick up the new dependencies."
        ) from exc

    if not _tesseract_binary_available(pytesseract):
        raise ExtractionError(
            "tesseract binary not found on PATH. Install it (macOS: "
            "`brew install tesseract`) and re-run the upload."
        )

    try:
        with Image.open(path) as img:
            # Force-load and convert to a mode tesseract handles cleanly.
            # Some PNGs come in palette mode and produce odd OCR output
            # without normalization.
            normalized = img.convert("RGB")
            text = pytesseract.image_to_string(normalized)
    except Exception as exc:
        # pytesseract raises ``TesseractError`` with the binary's stderr;
        # PIL raises ``UnidentifiedImageError`` for files we can't read.
        # Both are user-facing extraction failures.
        raise ExtractionError(f"OCR failed: {exc}") from exc

    return (text or "").strip(), "ocr"


def _tesseract_binary_available(pytesseract) -> bool:
    """Best-effort probe — does ``tesseract --version`` succeed?

    pytesseract caches the binary path in ``pytesseract.tesseract_cmd``
    and falls back to ``shutil.which("tesseract")``. We invoke
    ``get_tesseract_version`` because it actually shells out — the
    cheaper check (``shutil.which``) misses cases where the path is
    stale or the binary is broken.
    """
    try:
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False
